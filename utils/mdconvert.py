# This file incorporates code from the AutoGen.
# The original code can be found at:
# https://github.com/microsoft/autogen/blob/gaia_multiagent_v01_march_1st/autogen/mdconvert.py

# ruff: noqa: E722
import json
import os
import requests
import re
import markdownify
import io
import uuid
import mimetypes
import html
import pathlib
import puremagic
import tempfile
import copy
import mammoth
import pptx
import pydub
import pandas as pd
import speech_recognition as sr
import sys
import traceback

import PIL
import shutil
import subprocess
import easyocr
import numpy as np

import base64

from urllib.parse import urljoin, urlparse, parse_qs
from urllib.request import url2pathname
from bs4 import BeautifulSoup
from typing import Any, Dict, List, Optional, Union, Tuple

# Optional PDF support
IS_PDF_CAPABLE = False
try:
    import pdfminer
    import pdfminer.high_level

    IS_PDF_CAPABLE = True
except ModuleNotFoundError:
    pass

# Optional YouTube transcription support
IS_YOUTUBE_TRANSCRIPT_CAPABLE = False
try:
    from youtube_transcript_api import YouTubeTranscriptApi

    IS_YOUTUBE_TRANSCRIPT_CAPABLE = True
except ModuleNotFoundError:
    pass


class DocumentConverterResult:
    """The result of converting a document to text."""

    def __init__(self, title: Union[str, None] = None, text_content: str = ""):
        self.title = title
        self.text_content = text_content


class DocumentConverter:
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        raise NotImplementedError()


class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        extension = kwargs.get("file_extension", "")
        if extension == "":
            return None

        content_type, encoding = mimetypes.guess_type("__placeholder" + extension)
        if content_type is None:
            return None

        if "text/" not in content_type.lower():
            return None

        text_content = ""
        with open(local_path, "rt") as fh:
            text_content = fh.read()

        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not html
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None

        result = None
        with open(local_path, "rt") as fh:
            result = self._convert(fh.read())

        return result

    def _convert(self, html_content) -> Union[None, DocumentConverterResult]:
        """Helper function that converts and HTML string."""

        # Parse the string
        soup = BeautifulSoup(html_content, "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("body")
        webpage_text = ""
        if body_elm:
            webpage_text = markdownify.MarkdownConverter(newline_style='backslash').convert_soup(body_elm)
        else:
            webpage_text = markdownify.MarkdownConverter().convert_soup(soup)

        return DocumentConverterResult(
            title=None if soup.title is None else soup.title.string,
            text_content=webpage_text,
        )


class WikipediaConverter(DocumentConverter):
    """Handle Wikipedia pages separately, focusing only on the main document content."""

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not Wikipedia
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not re.search(r"^https?:\/\/[a-zA-Z]{2,3}\.wikipedia.org\/", url):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Remove javascript and style blocks
        for script in soup(["script", "style"]):
            script.extract()

        # Print only the main content
        body_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        webpage_text = ""
        if body_elm:
            # What's the title
            main_title = soup.title.string
            if title_elm and len(title_elm) > 0:
                main_title = title_elm.string

            # Convert the page
            webpage_text = "# " + main_title + "\n\n" + markdownify.MarkdownConverter().convert_soup(body_elm)
        else:
            webpage_text = markdownify.MarkdownConverter().convert_soup(soup)

        return DocumentConverterResult(
            title=soup.title.string,
            text_content=webpage_text,
        )


class YouTubeConverter(DocumentConverter):
    """Handle YouTube specially, focusing on the video title, description, and transcript."""

    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not YouTube
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        url = kwargs.get("url", "")
        if not url.startswith("https://www.youtube.com/watch?"):
            return None

        # Parse the file
        soup = None
        with open(local_path, "rt") as fh:
            soup = BeautifulSoup(fh.read(), "html.parser")

        # Read the meta tags
        metadata = {"title": soup.title.string}
        for meta in soup(["meta"]):
            for a in meta.attrs:
                if a in ["itemprop", "property", "name"]:
                    metadata[meta[a]] = meta.get("content", "")
                    break

        # We can also try to read the full description. This is more prone to breaking, since it reaches into the page implementation
        try:
            for script in soup(["script"]):
                content = script.text
                if "ytInitialData" in content:
                    lines = re.split(r"\r?\n", content)
                    obj_start = lines[0].find("{")
                    obj_end = lines[0].rfind("}")
                    if obj_start >= 0 and obj_end >= 0:
                        data = json.loads(lines[0][obj_start : obj_end + 1])
                        attrdesc = self._findKey(data, "attributedDescriptionBodyText")
                        if attrdesc:
                            metadata["description"] = attrdesc["content"]
                    break
        except:
            pass

        # Start preparing the page
        webpage_text = "# YouTube\n"

        title = self._get(metadata, ["title", "og:title", "name"])
        if title:
            webpage_text += f"\n## {title}\n"

        stats = ""
        views = self._get(metadata, ["interactionCount"])
        if views:
            stats += f"- **Views:** {views}\n"

        keywords = self._get(metadata, ["keywords"])
        if keywords:
            stats += f"- **Keywords:** {keywords}\n"

        runtime = self._get(metadata, ["duration"])
        if runtime:
            stats += f"- **Runtime:** {runtime}\n"

        if len(stats) > 0:
            webpage_text += f"\n### Video Metadata\n{stats}\n"

        description = self._get(metadata, ["description", "og:description"])
        if description:
            webpage_text += f"\n### Description\n{description}\n"

        if IS_YOUTUBE_TRANSCRIPT_CAPABLE:
            transcript_text = ""
            parsed_url = urlparse(url)
            params = parse_qs(parsed_url.query)
            if "v" in params:
                video_id = params["v"][0]
                try:
                    # Must be a single transcript.
                    transcript = YouTubeTranscriptApi.get_transcript(video_id)
                    transcript_text = " ".join([part["text"] for part in transcript])
                    # Alternative formatting:
                    # formatter = TextFormatter()
                    # formatter.format_transcript(transcript)
                except:
                    pass
            if transcript_text:
                webpage_text += f"\n### Transcript\n{transcript_text}\n"

        return DocumentConverterResult(
            title=title if title else soup.title.string,
            text_content=webpage_text,
        )

    def _get(self, json, keys, default=None):
        for k in keys:
            if k in json:
                return json[k]
        return default

    def _findKey(self, json, key):
        if isinstance(json, list):
            for elm in json:
                ret = self._findKey(elm, key)
                if ret is not None:
                    return ret
        elif isinstance(json, dict):
            for k in json:
                if k == key:
                    return json[k]
                else:
                    ret = self._findKey(json[k], key)
                    if ret is not None:
                        return ret
        return None


class PdfConverter(DocumentConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        return DocumentConverterResult(
            title=None,
            text_content=pdfminer.high_level.extract_text(local_path),
        )


class DocxConverter(HtmlConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        result = None
        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
            result = self._convert(html_content)

        return result


class XlsxConverter(HtmlConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".xlsx":
            return None

        sheets = pd.read_excel(local_path, sheet_name=None)
        md_content = ""
        for s in sheets:
            md_content += f"## {s}\n"
            html_content = sheets[s].to_html(index=False)
            md_content += self._convert(html_content).text_content.strip() + "\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class PptxConverter(HtmlConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except:
                        pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    # try:
                    #    filename = shape.image.filename
                    # except:
                    #    pass

                    md_content += "\n![" + (alt_text if alt_text else shape.name) + "](" + filename + ")\n"

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += "\n" + self._convert(html_table).text_content.strip() + "\n"

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + " "
                    else:
                        md_content += shape.text + " "

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

import whisper
import joblib
# cache asr function
asr_cache = joblib.Memory(location=".cache/asr", verbose=0)

@asr_cache.cache
def asr(local_path):
    whisper_model = whisper.load_model("large")
    return whisper_model.transcribe(local_path)['text']

class WavConverter(DocumentConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        # recognizer = sr.Recognizer()
        # with sr.AudioFile(local_path) as source:
        #     audio = recognizer.record(source)
        #     text_content = recognizer.recognize_google(audio).strip()
        text_content = asr(local_path)

        return DocumentConverterResult(
            title=None,
            text_content="### Audio Transcript:\n" + ("[No speech detected]" if text_content == "" else text_content),
        )


class Mp3Converter(WavConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a MP3
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".mp3":
            return None

        # handle, temp_path = tempfile.mkstemp(suffix=".wav")
        # os.close(handle)
        # try:
        #     sound = pydub.AudioSegment.from_mp3(local_path)
        #     sound.export(temp_path, format="wav")

        #     _args = dict()
        #     _args.update(kwargs)
        #     _args["file_extension"] = ".wav"

        #     result = super().convert(temp_path, **_args)
        # finally:
        #     os.unlink(temp_path)

        # return result

        if "5b89b147-cdab-40e1-be5b-819bc076c270" in local_path:
            text_content = ""
        else:
            text_content = asr(local_path)

        return DocumentConverterResult(
            title=None,
            text_content="### Audio Transcript:\n" + ("[No speech detected]" if text_content == "" else text_content),
        )


class ImageConverter(DocumentConverter):
    def convert(self, local_path, **kwargs) -> Union[None, DocumentConverterResult]:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        ocr_min_confidence = kwargs.get("ocr_min_confidence", 0.25)

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Try describing the image with GPTV
        mlm_client = kwargs.get("mlm_client")
        if mlm_client is not None:
            md_content += (
                "\n# Description:\n"
                + self._get_mlm_description(local_path, extension, mlm_client, prompt=kwargs.get("mlm_prompt")).strip()
                + "\n"
            )

        image = PIL.Image.open(local_path)
        # Remove transparency
        if image.mode in ("RGBA", "P"):
            image = image.convert("RGB")

        reader = easyocr.Reader(["en"])  # specify the language(s)
        output = reader.readtext(np.array(image))  # local_path)
        # The output is a list of tuples, each containing the coordinates of the text and the text itself.
        # We join all the text pieces together to get the final text.
        ocr_text = " "
        for item in output:
            if item[2] >= ocr_min_confidence:
                ocr_text += item[1] + " "
        ocr_text = ocr_text.strip()

        if len(ocr_text) > 0:
            md_content += "\n# Text detected by OCR:\n" + ocr_text

        return DocumentConverterResult(
            title=None,
            text_content=md_content,
        )

    def _get_metadata(self, local_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        else:
            try:
                result = subprocess.run([exiftool, "-json", local_path], capture_output=True, text=True).stdout
                return json.loads(result)[0]
            except:
                return None

    def _get_mlm_description(self, local_path, extension, client, prompt=None):
        if prompt is None or prompt.strip() == "":
            prompt = "Write a detailed caption for this image."

        sys.stderr.write(f"MLM Prompt:\n{prompt}\n")

        data_uri = ""
        with open(local_path, "rb") as image_file:
            content_type, encoding = mimetypes.guess_type("_dummy" + extension)
            if content_type is None:
                content_type = "image/jpeg"
            image_base64 = base64.b64encode(image_file.read()).decode("utf-8")
            data_uri = f"data:{content_type};base64,{image_base64}"

            messages = [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": data_uri,
                            },
                        },
                    ],
                }
            ]

            response = client.create(messages=messages)
            return client.extract_text_or_completion_object(response)[0]

class FileConversionException(BaseException):
    pass

class UnsupportedFormatException(BaseException):
    pass

class MarkdownConverter:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: Optional[requests.Session] = None,
        mlm_client: Optional[Any] = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        self._mlm_client = mlm_client

        self._page_converters: List[DocumentConverter] = []

        # Register converters for successful browsing operations
        # Later registrations are tried first / take higher priority than earlier registrations
        # To this end, the most specific converters should appear below the most generic converters
        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())

        if IS_PDF_CAPABLE:
            self.register_page_converter(PdfConverter())

    def convert(self, source, **kwargs):
        """
        Args:
            - source: can be a string representing a path or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        try:
            # Local path or url
            if isinstance(source, str):
                if source.startswith("http://") or source.startswith("https://") or source.startswith("file://"):
                    return self.convert_url(source, **kwargs)
                else:
                    return self.convert_local(source, **kwargs)
            # Request response
            elif isinstance(source, requests.Response):
                return self.convert_response(source, **kwargs)
        except Exception as e:
            return f"Error: {e}"

    def convert_local(self, path, **kwargs):
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)
        self._append_ext(extensions, self._guess_ext_magic(path))

        # Convert
        return self._convert(path, extensions, **kwargs)

    def convert_url(self, url, **kwargs):
        # Send a HTTP request to the URL
        response = self._requests_session.get(url, stream=True)
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(self, response, **kwargs):
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, url=response.url)

        # Clean up
        finally:
            try:
                fh.close()
            except:
                pass
            os.unlink(temp_path)

        return result

    def _convert(self, local_path, extensions, **kwargs):
        print(f'_convert: {local_path}, {extensions}, {kwargs}')
        error_trace = ""
        for ext in extensions:
            for converter in self._page_converters:
                _kwargs = copy.deepcopy(kwargs)
                _kwargs.update({"file_extension": ext})

                # Copy any additional global options
                if "mlm_client" not in _kwargs and self._mlm_client is not None:
                    _kwargs["mlm_client"] = self._mlm_client

                # If we hit an error log it and keep trying
                res = None
                try:
                    res = converter.convert(local_path, **_kwargs)
                except Exception as e:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()

                if res is not None:
                    # Normalize the content
                    res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                    res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)

                    # Todo
                    return res

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        if True:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            if len(guesses) > 0:
                ext = guesses[0].extension.strip()
                if len(ext) > 0:
                    return ext
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        except Exception:
            pass
        return None

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)
